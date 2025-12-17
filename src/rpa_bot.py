from supabase import create_client
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
from google.api_core import exceptions as google_exceptions
import pandas as pd
import os


load_dotenv()

# Function to get data from Supabase
def get_data_from_supabase():
    supabase_url = os.getenv("SUPABASE_URL")
    supabase_key = os.getenv("SUPABASE_KEY")
    supabase = create_client(supabase_url, supabase_key)

    response = supabase.schema("public").table('financial_data').select("*").execute()
    df = pd.DataFrame(response.data)
    return df

# Function to get data from local Excel file
def get_data_from_local(file_path):
    df = pd.read_excel(file_path)
    return df

# Function to prepare data for charts
def prepare_data_for_chart(df):

    # Chart 1: Pie Chart - Spending by Category
    pie_data = df.groupby('category')['amount'].sum().reset_index()
    pie_data = pie_data.sort_values('amount', ascending=False)
    pie_data['percentage'] = (pie_data['amount'] / pie_data['amount'].sum() * 100).round(2)
    
    # Chart 2: Bar Chart - Quarterly Spending by Category
    bar_data = df.groupby(['quarter', 'category'])['amount'].sum().reset_index()
    
    # Chart 3: Line Chart - Monthly Spending Trend
    df['month'] = pd.to_datetime(df['transaction_date']).dt.to_period('M').astype(str)
    line_data = df.groupby('month')['amount'].sum().reset_index()
    line_data = line_data.sort_values('month')
    
    # Chart 4: Stacked Bar Chart - Category Breakdown by Quarter
    stacked_data = df.groupby(['quarter', 'category'])['amount'].sum().reset_index()

    return pie_data, bar_data, line_data, stacked_data

# Function to prompt GenAI for analysis
@retry(
    retry=retry_if_exception_type(google_exceptions.ServerError),
    wait=wait_exponential(multiplier=1, min=4, max=10),
    stop=stop_after_attempt(3)
)
def promt_to_genai(df):
    pie_data, bar_data, line_data, stacked_data = prepare_data_for_chart(df)

    # Prepare data frame for GenAI prompt
    pie_chart_data_summary = f"""
    Pie Chart Data (Total Spending by Category):
    {pie_data.to_string(index=False)}

    Total Spending: ${df['amount'].sum():.2f}
    Date range: {df['transaction_date'].min()} to {df['transaction_date'].max()}
    """

    bar_chart_data_summary = f"""
    Bar Chart Data (Quarterly Spending Trends):
    {bar_data.to_string(index=False)}

    Total Spending: ${df['amount'].sum():.2f}
    Date range: {df['transaction_date'].min()} to {df['transaction_date'].max()}
    """

    # Initialize GenAI client
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    # Set up the model
    model = genai.GenerativeModel(model_name="gemini-2.5-flash")

    # Generate content using GenAI
    prompt_summarize_pie_chart = f"""Analyze this financial data and summarize in 50-100 words your analysis to put into presentation slides.

                    {pie_chart_data_summary}

                    Please provide: Key insights from the spending distribution"""

    prompt_summarize_bar_chart = f"""Analyze this financial data and summarize in 50-100 words your analysis to put into presentation slides.

                    {bar_chart_data_summary}

                    Please provide: Notable trends across quarters"""
    
    response = {}

    response['pie_description'] = model.generate_content(prompt_summarize_pie_chart).text
    response['bar_description'] = model.generate_content(prompt_summarize_bar_chart).text

    return response

# Function to update charts in a PowerPoint presentation
def update_presentation(pie_data, bar_data, line_data, stacked_data, template_path, output_path, descriptions=None):
    """
    Updates charts in a PowerPoint presentation from a template.
    - Slide 1 is expected to have a Pie Chart.
    - Slide 2 is expected to have a Clustered Column Chart.
    """
    prs = Presentation(template_path)

    # --- Update Pie Chart on Slide 1 ---
    slide_0 = prs.slides[0]
    # Assume the chart is the first shape of type 'chart' on the slide
    pie_chart_shape = [s for s in slide_0.shapes if s.has_chart][0]
    pie_chart = pie_chart_shape.chart

    chart_data = CategoryChartData()
    chart_data.categories = pie_data['category']
    chart_data.add_series('Amount', pie_data['amount'])
    pie_chart.replace_data(chart_data)

    # Set chart title for Pie Chart
    if pie_chart.has_title:
        pie_chart.chart_title.text_frame.text = "Spending by Category"




    bar_chart_shape = [s for s in slide_0.shapes if s.has_chart][1]
    bar_chart = bar_chart_shape.chart

    # Pivot data to have quarters as categories and spending categories as series
    pivoted_bar_data = bar_data.pivot(index='quarter', columns='category', values='amount').fillna(0)

    chart_data = CategoryChartData()
    chart_data.categories = pivoted_bar_data.index.astype(str) # Quarters

    for col in pivoted_bar_data.columns:
        chart_data.add_series(col, pivoted_bar_data[col]) # Add each spending category as a series

    bar_chart.replace_data(chart_data)
    
    # Set chart title for Bar Chart
    if bar_chart.has_title:
        bar_chart.chart_title.text_frame.text = "Quarterly by Category"

    # Update text description for Pie Chart
    if descriptions and 'pie_description' in descriptions:
        for shape in slide_0.shapes:
            if shape.has_text_frame and not shape.has_chart and len(shape.text_frame.text) > 30:
                shape.text_frame.clear()  # Clear existing text
                text_frame = shape.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = descriptions['pie_description']
                for run in p.runs:
                    run.font.name = 'Century Gothic'
                break


    # --- Update Line Chart on Slide 2 ---
    slide_1 = prs.slides[1]
    bar_chart_shape = [s for s in slide_1.shapes if s.has_chart][0]
    bar_chart = bar_chart_shape.chart

    # Pivot data to have quarters as categories and spending categories as series
    pivoted_bar_data = bar_data.pivot(index='quarter', columns='category', values='amount').fillna(0)

    chart_data = CategoryChartData()
    chart_data.categories = pivoted_bar_data.index.astype(str) # Quarters

    for col in pivoted_bar_data.columns:
        chart_data.add_series(col, pivoted_bar_data[col]) # Add each spending category as a series

    bar_chart.replace_data(chart_data)
    
    # Set chart title for Bar Chart
    if bar_chart.has_title:
        bar_chart.chart_title.text_frame.text = "Quarterly by Category"


    # Assume the chart is the first shape of type 'chart' on the slide
    line_chart_shape = [s for s in slide_1.shapes if s.has_chart][1]
    line_chart = line_chart_shape.chart

    # Line data is already aggregated by month
    chart_data = CategoryChartData()
    chart_data.categories = line_data['month'].astype(str)  # Months as categories
    chart_data.add_series('Monthly Spending', line_data['amount'])

    line_chart.replace_data(chart_data)
    
    # Set chart title for Line Chart
    if line_chart.has_title:
        line_chart.chart_title.text_frame.text = "Monthly Spending Trend"

    

    # Update text description for Bar Chart
    if descriptions and 'bar_description' in descriptions:
        for shape in slide_1.shapes:
            if shape.has_text_frame and not shape.has_chart and len(shape.text_frame.text) > 30:
                shape.text_frame.clear()
                text_frame = shape.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = descriptions['bar_description']
                for run in p.runs:
                    run.font.name = 'Century Gothic'
                break

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    # Construct the path to the data file relative to the script's location
    script_dir = os.path.dirname(__file__)
    file_path = os.path.join(script_dir, '..', 'data', 'xlsx', 'financial_data.xlsx')
    # df = get_data_from_local(file_path)
    df = get_data_from_supabase()
    pie_data, bar_data, line_data, stacked_data = prepare_data_for_chart(df)

    # Define presentation paths
    template_presentation_path = os.path.join(script_dir, '..', 'data', 'power-point', 'template.pptx')
    output_presentation_path = os.path.join(script_dir, '..', 'data', 'power-point', 'financial_report.pptx')

    # Update the presentation with new data
    analysis = promt_to_genai(df)
    update_presentation(pie_data, bar_data, line_data, stacked_data, template_presentation_path, output_presentation_path, analysis)